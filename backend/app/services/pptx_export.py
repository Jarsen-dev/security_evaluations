"""Exportación del reporte a PowerPoint 16:9, lista para presentar a gerencia.

Las gráficas son NATIVAS de PowerPoint (``chart_data`` + ``add_chart``), no
imágenes: quien presente puede editarlas, cambiar colores o corregir un dato
sin volver a generar el archivo.
"""

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from app.core.config import RUTA_LOGO, hay_logo
from app.services.exportacion_comun import (
    DatosReporte,
    generar_conclusiones,
    preguntas_mas_falladas,
)

# 16:9
ANCHO = Inches(13.333)
ALTO = Inches(7.5)

AZUL = RGBColor(0x1F, 0x4E, 0x79)
AZUL_CLARO = RGBColor(0x2F, 0x81, 0xF7)
GRIS_TEXTO = RGBColor(0x59, 0x59, 0x59)
VERDE = RGBColor(0x0F, 0x7B, 0x2F)
ROJO = RGBColor(0xC0, 0x26, 0x26)

MAX_TEXTO_PREGUNTA = 55


def _diapositiva_en_blanco(presentacion: Presentation) -> Slide:
    """Crea una diapositiva sin marcadores de posición."""
    return presentacion.slides.add_slide(presentacion.slide_layouts[6])


def _titulo(diapositiva: Slide, texto: str, subtitulo: str | None = None) -> None:
    """Escribe el título estándar de una diapositiva de contenido."""
    caja = diapositiva.shapes.add_textbox(
        Inches(0.6), Inches(0.35), ANCHO - Inches(1.2), Inches(0.9)
    )
    marco = caja.text_frame
    marco.word_wrap = True

    parrafo = marco.paragraphs[0]
    parrafo.text = texto
    parrafo.font.size = Pt(28)
    parrafo.font.bold = True
    parrafo.font.color.rgb = AZUL

    if subtitulo:
        linea = marco.add_paragraph()
        linea.text = subtitulo
        linea.font.size = Pt(13)
        linea.font.color.rgb = GRIS_TEXTO


def _recortar(texto: str, limite: int = MAX_TEXTO_PREGUNTA) -> str:
    """Acorta el texto de una pregunta para que quepa en el eje."""
    return texto if len(texto) <= limite else f"{texto[: limite - 1]}…"


def _estilo_grafica(grafica, con_leyenda: bool = False) -> None:
    """Aplica el estilo común a las gráficas nativas."""
    grafica.has_title = False
    grafica.font.size = Pt(12)

    if con_leyenda:
        grafica.has_legend = True
        grafica.legend.position = XL_LEGEND_POSITION.BOTTOM
        grafica.legend.include_in_layout = False
    else:
        grafica.has_legend = False


# --- 1. Portada ------------------------------------------------------------


ANCHO_LOGO_PORTADA = Inches(2.1)


def _portada(presentacion: Presentation, datos: DatosReporte, periodo: str) -> None:
    """Logo a la izquierda y el título a su derecha, sobre fondo claro.

    El logo institucional es azul marino: sobre una franja del mismo tono
    quedaría invisible. Por eso la portada usa fondo blanco con una barra de
    acento delgada arriba, en lugar del bloque de color completo.
    """
    diapositiva = _diapositiva_en_blanco(presentacion)

    barra = diapositiva.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, Inches(0.28))
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL
    barra.line.fill.background()
    barra.shadow.inherit = False

    izquierda_texto = Inches(1.0)
    arriba_bloque = Inches(2.3)

    if hay_logo():
        imagen = diapositiva.shapes.add_picture(
            str(RUTA_LOGO), Inches(1.0), arriba_bloque, width=ANCHO_LOGO_PORTADA
        )
        # El título arranca después del logo, con un respiro de media pulgada.
        izquierda_texto = Emu(int(imagen.left + imagen.width + Inches(0.5)))
        # Se centra verticalmente el bloque de texto respecto al logo.
        arriba_texto = Emu(int(imagen.top - Inches(0.15)))
    else:
        arriba_texto = arriba_bloque

    caja = diapositiva.shapes.add_textbox(
        izquierda_texto,
        arriba_texto,
        ANCHO - izquierda_texto - Inches(0.9),
        Inches(2.0),
    )
    marco = caja.text_frame
    marco.word_wrap = True

    parrafo = marco.paragraphs[0]
    parrafo.text = "Evaluación de Conocimientos"
    parrafo.font.size = Pt(16)
    parrafo.font.color.rgb = GRIS_TEXTO

    nombre = marco.add_paragraph()
    nombre.text = datos.cuestionario.nombre
    nombre.font.size = Pt(32)
    nombre.font.bold = True
    nombre.font.color.rgb = AZUL

    detalle = diapositiva.shapes.add_textbox(
        Inches(1.0), Inches(5.0), ANCHO - Inches(2.0), Inches(1.8)
    )
    marco_detalle = detalle.text_frame
    marco_detalle.word_wrap = True

    for indice, linea in enumerate(
        [
            f"Periodo: {periodo}",
            f"Respuestas recibidas: {datos.resumen['total_respuestas']}",
            f"Generado el {datos.generado_at.strftime('%d/%m/%Y %H:%M')} UTC",
        ]
    ):
        parrafo_detalle = (
            marco_detalle.paragraphs[0] if indice == 0 else marco_detalle.add_paragraph()
        )
        parrafo_detalle.text = linea
        parrafo_detalle.font.size = Pt(15)
        parrafo_detalle.font.color.rgb = GRIS_TEXTO


# --- 2. Resumen ejecutivo --------------------------------------------------


def _resumen_ejecutivo(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    _titulo(diapositiva, "Resumen ejecutivo")

    resumen = datos.resumen
    participacion = resumen["participacion"]

    tarjetas = [
        ("Respuestas recibidas", str(resumen["total_respuestas"]), AZUL_CLARO),
        (
            "Nivel de participación",
            f"{participacion['porcentaje']:.1f}%"
            if participacion["porcentaje"] is not None
            else "Sin meta",
            AZUL_CLARO,
        ),
        (
            "Calificación promedio",
            f"{resumen['promedio_general']:.1f}%"
            if resumen["promedio_general"] is not None
            else "—",
            VERDE
            if (resumen["promedio_general"] or 0) >= resumen["umbral_aprobacion"]
            else ROJO,
        ),
        (
            "Tasa de aprobación",
            f"{resumen['tasa_aprobacion']:.1f}%"
            if resumen["tasa_aprobacion"] is not None
            else "—",
            VERDE
            if (resumen["tasa_aprobacion"] or 0) >= resumen["umbral_aprobacion"]
            else ROJO,
        ),
    ]

    ancho_tarjeta = Inches(2.9)
    separacion = Inches(0.25)
    inicio = Inches(0.6)
    arriba = Inches(2.1)

    for indice, (etiqueta, valor, color) in enumerate(tarjetas):
        izquierda = Emu(int(inicio + indice * (ancho_tarjeta + separacion)))

        caja = diapositiva.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, izquierda, arriba, ancho_tarjeta, Inches(2.4)
        )
        caja.fill.solid()
        caja.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        caja.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
        caja.shadow.inherit = False

        marco = caja.text_frame
        marco.word_wrap = True
        marco.margin_top = Inches(0.3)

        parrafo_valor = marco.paragraphs[0]
        parrafo_valor.text = valor
        parrafo_valor.alignment = PP_ALIGN.CENTER
        parrafo_valor.font.size = Pt(40)
        parrafo_valor.font.bold = True
        parrafo_valor.font.color.rgb = color

        parrafo_etiqueta = marco.add_paragraph()
        parrafo_etiqueta.text = etiqueta
        parrafo_etiqueta.alignment = PP_ALIGN.CENTER
        parrafo_etiqueta.font.size = Pt(14)
        parrafo_etiqueta.font.color.rgb = GRIS_TEXTO

    pie = diapositiva.shapes.add_textbox(
        Inches(0.6), Inches(5.0), ANCHO - Inches(1.2), Inches(0.6)
    )
    pie.text_frame.paragraphs[0].text = (
        f"Umbral de aprobación: {resumen['umbral_aprobacion']}%  ·  "
        f"{resumen['aprobados']} aprobados de {resumen['total_respuestas']}  ·  "
        f"{resumen['total_en_progreso']} intentos sin finalizar"
    )
    pie.text_frame.paragraphs[0].font.size = Pt(13)
    pie.text_frame.paragraphs[0].font.color.rgb = GRIS_TEXTO


# --- 3. Participación por área --------------------------------------------


def _participacion(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    _titulo(
        diapositiva,
        "Participación por área",
        "Respuestas recibidas contra la meta de headcount.",
    )

    areas = datos.por_area
    grafica_datos = CategoryChartData()
    grafica_datos.categories = [area["label"] for area in areas]
    grafica_datos.add_series("Respuestas", [area["intentos"] for area in areas])

    if any(area["meta"] for area in areas):
        # Las áreas sin meta van en cero para no romper la serie; su columna
        # simplemente no aparece.
        grafica_datos.add_series("Meta", [area["meta"] or 0 for area in areas])

    marco = diapositiva.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(1.6),
        ANCHO - Inches(1.2),
        Inches(5.2),
        grafica_datos,
    )
    _estilo_grafica(marco.chart, con_leyenda=True)


# --- 4. Calificación promedio por área ------------------------------------


def _promedio_por_area(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    umbral = datos.resumen["umbral_aprobacion"]
    _titulo(
        diapositiva,
        "Calificación promedio por área",
        f"Ordenadas de mayor a menor. Umbral de aprobación: {umbral}%.",
    )

    con_datos = sorted(
        [area for area in datos.por_area if area["promedio"] is not None],
        key=lambda area: area["promedio"] or 0,
    )

    grafica_datos = CategoryChartData()
    grafica_datos.categories = [area["label"] for area in con_datos] or ["Sin datos"]
    grafica_datos.add_series(
        "Promedio %", [area["promedio"] for area in con_datos] or [0]
    )

    marco = diapositiva.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.6),
        Inches(1.6),
        ANCHO - Inches(1.2),
        Inches(5.2),
        grafica_datos,
    )
    _estilo_grafica(marco.chart)

    grafica = marco.chart
    grafica.value_axis.maximum_scale = 100
    grafica.plots[0].has_data_labels = True
    grafica.plots[0].data_labels.number_format = "0.0"
    grafica.plots[0].data_labels.number_format_is_linked = False


# --- 5. Distribución de calificaciones ------------------------------------


def _distribucion(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    _titulo(
        diapositiva,
        "Distribución de calificaciones",
        "Cuántas personas cayeron en cada rango.",
    )

    grafica_datos = CategoryChartData()
    grafica_datos.categories = [rango["rango"] for rango in datos.distribucion]
    grafica_datos.add_series(
        "Personas", [rango["cantidad"] for rango in datos.distribucion]
    )

    marco = diapositiva.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(1.6),
        ANCHO - Inches(1.2),
        Inches(5.2),
        grafica_datos,
    )
    _estilo_grafica(marco.chart)

    marco.chart.plots[0].has_data_labels = True


# --- 6. Preguntas con mayor error -----------------------------------------


def _preguntas_falladas(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    _titulo(
        diapositiva,
        "Top 10 preguntas con mayor índice de error",
        "Señala qué temas necesitan recapacitación o qué preguntas están mal redactadas.",
    )

    peores = preguntas_mas_falladas(datos.por_pregunta)

    grafica_datos = CategoryChartData()
    # Se invierte el orden: en una barra horizontal, la primera categoría se
    # dibuja abajo, y la peor pregunta debe quedar arriba.
    grafica_datos.categories = [
        _recortar(pregunta["texto"]) for pregunta in reversed(peores)
    ] or ["Sin datos"]
    grafica_datos.add_series(
        "% de error",
        [pregunta["porcentaje_error"] for pregunta in reversed(peores)] or [0],
    )

    marco = diapositiva.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.6),
        Inches(1.7),
        ANCHO - Inches(1.2),
        Inches(5.1),
        grafica_datos,
    )
    _estilo_grafica(marco.chart)

    grafica = marco.chart
    grafica.value_axis.maximum_scale = 100
    grafica.font.size = Pt(10)
    grafica.plots[0].has_data_labels = True
    grafica.plots[0].data_labels.number_format = "0.0"
    grafica.plots[0].data_labels.number_format_is_linked = False

    for serie in grafica.series:
        serie.format.fill.solid()
        serie.format.fill.fore_color.rgb = ROJO


# --- 7. Conclusiones -------------------------------------------------------


def _conclusiones(presentacion: Presentation, datos: DatosReporte) -> None:
    diapositiva = _diapositiva_en_blanco(presentacion)
    _titulo(diapositiva, "Conclusiones")

    caja = diapositiva.shapes.add_textbox(
        Inches(0.8), Inches(1.7), ANCHO - Inches(1.6), Inches(5.2)
    )
    marco = caja.text_frame
    marco.word_wrap = True

    for indice, linea in enumerate(generar_conclusiones(datos)):
        parrafo = marco.paragraphs[0] if indice == 0 else marco.add_paragraph()
        # Las líneas numeradas ya vienen indentadas desde el generador.
        es_sublinea = linea.startswith("    ")
        parrafo.text = linea if es_sublinea else f"•  {linea}"
        parrafo.font.size = Pt(14 if es_sublinea else 16)
        parrafo.font.color.rgb = GRIS_TEXTO if es_sublinea else RGBColor(0x22, 0x22, 0x22)
        parrafo.space_after = Pt(10)


def generar_pptx(datos: DatosReporte, periodo: str) -> BytesIO:
    """Arma la presentación de siete diapositivas y la devuelve en memoria."""
    presentacion = Presentation()
    presentacion.slide_width = ANCHO
    presentacion.slide_height = ALTO

    _portada(presentacion, datos, periodo)
    _resumen_ejecutivo(presentacion, datos)
    _participacion(presentacion, datos)
    _promedio_por_area(presentacion, datos)
    _distribucion(presentacion, datos)
    _preguntas_falladas(presentacion, datos)
    _conclusiones(presentacion, datos)

    flujo = BytesIO()
    presentacion.save(flujo)
    flujo.seek(0)

    return flujo
